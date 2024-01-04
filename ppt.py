import streamlit as st
from pptx import Presentation

def ppt_converter(prs):
    if prs is not None:
        presentation = Presentation(prs)

        # Extract text from all slides (both title and body text)
        all_text = ""
        for slide in presentation.slides:
            title_text = slide.shapes.title.text if slide.shapes.title else ""
            body_text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, 'text')])
            
            slide_text = f"Slide {slide.slide_id-255}: {title_text}\n{body_text}\n"
            all_text += slide_text

        return all_text
